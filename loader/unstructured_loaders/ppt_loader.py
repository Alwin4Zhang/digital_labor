# -*- coding: utf-8 -*-
"""
  @CreateTime	:  2024/01/11 17:43:43
  @Author	:  Alwin Zhang
  @Mail	:  zhangjunfeng@rainbowcn.com
"""

import os
import sys

sys.path.append("./")

import uuid
import re
import base64
import logging
from copy import deepcopy
from typing import List, Tuple, Dict, Optional, Any
from io import BytesIO
from pprint import pprint

from tqdm import tqdm
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
# from rapidocr_onnxruntime import RapidOCR

from loader.unstructured_loaders.constants import *
from loader.unstructured_loaders.loader import FileLoader
from utils.common import convert_base64_to_buffer, upload_file, pt2px,ocr_api
from utils.string_util import normalize
from configs.apollo_config import logger

# ocr = RapidOCR()



class PPTLoader(FileLoader):

    def __init__(self, file_path: str = None, file_stream: BytesIO = None):
        super().__init__(file_path, file_stream)
        self.images = {}
        self.blocks, self.maybe_title = self.ppt2blocks()
        self.blocks = self.block2html()
        self.tree = None

    def table2pd(self, table):
        df = [["" for i in range(len(table.columns))] for j in range(len(table.rows))]
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                if cell.text:
                    df[i][j] = normalize(cell.text).replace("\n", "<br>")
        if len(df) <= 1:
            return pd.DataFrame(columns=df[0])
        return pd.DataFrame(index=None).from_records(df[1:], columns=df[0])

    def ppt2blocks(self) -> Tuple:
        """ppt -> slide(contain shapes)"""
        doc = None
        if self.file_path is not None and os.path.exists(self.file_path):
            doc = Presentation(self.file_path)
        elif self.file_stream is not None:
            doc = Presentation(self.file_stream)
        blocks = []
        for i, slide in tqdm(enumerate(doc.slides)):
            page_title = None
            title = slide.shapes.title
            if title:
                page_title = normalize(title.text)
                # title_top, title_left, title_width, title_height = title.top.mm, title.left.mm, title.width.mm, title.height.mm
            shapes = []
            sorted_shapes = sorted(slide.shapes, key=lambda x: (x.top, x.left))
            for j, shape in enumerate(sorted_shapes):
                # top, left, width, height = shape.top.mm, shape.left.mm, shape.width.mm, shape.height.mm
                top, left, width, height = (
                    pt2px(shape.top.pt),
                    pt2px(shape.left.pt),
                    pt2px(shape.width.pt),
                    pt2px(shape.height.pt),
                )
                shape_id, shape_type = shape.shape_id, shape.shape_type
                d = {
                    "page_index": i,
                    "shape_index": j,
                    "top": top,
                    "left": left,
                    "width": width,
                    "height": height,
                    "name": shape.name,
                }
                if shape.has_text_frame:  # 文本框
                    d.update(
                        {
                            "type": "text",
                            "style": "Normal",
                            "data": normalize(shape.text),
                        }
                    )
                if shape_type == MSO_SHAPE_TYPE.GROUP: # GroupShape
                    text_group = ""
                    for shape in shape.shapes:
                        if shape.has_text_frame:
                            text_group += shape.text + "\n"
                    d.update(
                        {
                            "type": "text",
                            "style": "Normal",
                            "data": text_group,
                        }
                    )
                if shape_type == MSO_SHAPE_TYPE.TABLE:  # 表格
                    tb = self.table2pd(shape.table)
                    d.update({"type": "table", "style": "table", "data": tb})
                if shape_type == MSO_SHAPE_TYPE.CHART:  # 图表，eg.折线图、饼图等
                    # TODO: 后期解析图表chart
                    continue
                if shape_type == MSO_SHAPE_TYPE.PICTURE:  # 图片
                    image = shape.image
                    size = image.size  # 分辨率
                    filename = image.filename
                    uuid1 = uuid.uuid1().hex
                    image_base64 = base64.b64encode(image._blob)
                    ocr_text = ocr_api(image_base64)
                    ext = image.ext
                    content_type = image.content_type
                    d.update(
                        {
                            "type": "image",
                            "style": "image",
                            "size": size,
                            "filename": filename,
                            "uuid1": uuid1,
                            "ext": ext,
                            "content_type": content_type,
                            "data": image_base64,
                            "text": ocr_text
                        }
                    )

                if "data" in d:
                    shapes.append(d)
            blocks.append({"title": page_title, "shapes": shapes})
        mb_title = blocks[0]["title"]
        file_name = self._validate_filename()
        if file_name:
            mb_title = file_name
        return blocks, mb_title

    def block2html(self):
        # html渲染固定开头
        HTML_PREFIX = """<div id="page{index}">"""
        # html渲染固定结尾
        HTML_SUFFIX = """</div>"""
        blocks = []
        for i, slide in enumerate(self.blocks):
            this_html = HTML_PREFIX.format(index=i)
            this_markdown = ""
            page_title = slide["title"]
            shapes = slide["shapes"]
            raw_text = ""
            for j, shape in enumerate(shapes):
                style = shape["style"]
                if style in self.STYLE_MAP:
                    text = shape["data"]
                    if style == "image":
                        width = shape.get("width", 100)
                        height = shape.get("height", 100)
                        file_name = shape.get("filename", uuid.uuid1().hex + ".png")
                        buffer_stream = convert_base64_to_buffer(b64_data=text)
                        if buffer_stream in self.images:
                            assert_url = self.images.get(buffer_stream,"")
                        else:
                            assert_url = upload_file(
                                file_name=file_name, file_stream=buffer_stream
                            )
                            self.images[buffer_stream] = assert_url
                        if not assert_url:
                            continue
                        block_html = self.text2html(
                            assert_url, style, width=width, height=height
                        )
                        block_markdown = self.text2markdown(
                            text=assert_url, style=style
                        )
                    else:
                        raw_text += text + "\n"
                        block_html = self.text2html(text, style)
                        block_markdown = self.text2markdown(text=text, style=style)
                elif style == "table":
                    table = shape["data"]
                    table.style.set_table_styles(
                        [
                            {
                                "selector": "th,td",
                                "props": [
                                    ("font-size", "12pt"),
                                    ("border-style", "solid"),
                                    ("border-width", "1px"),
                                ],
                            }
                        ]
                    )
                    block_html = table.to_html(
                        index=False,
                        justify="left",
                        classes="mystyle",
                        escape=False,  # convert "< > &" to html style
                    )
                    block_markdown = table.to_markdown(index=False)
                this_html += block_html
                this_markdown += block_markdown
            this_html += HTML_SUFFIX
            blocks.append(
                {
                    "shapes": shapes,
                    "title": page_title,
                    "html": this_html,
                    "markdown": this_markdown,
                    "raw_text": raw_text,
                }
            )
        return blocks
